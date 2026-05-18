from pptx import Presentation
import os

class PPTGenerator:
    def __init__(self, manifest):
        self.manifest = manifest
        
    def generate(self, template_name, content_map, output_name):
        template_path = os.path.join(self.manifest['paths']['templates'], template_name)
        prs = Presentation(template_path)
        
        # Logic: Always clone from Slide Layouts (Indices) to preserve Master backgrounds
        # We assume index 1 is a standard content layout for this example
        layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(layout)
        
        for shape in slide.placeholders:
            # We map content based on placeholder name or index
            if shape.name in content_map:
                shape.text = content_map[shape.name]
                
        output_path = os.path.join(self.manifest['paths']['outputs'], f"{output_name}.pptx")
        prs.save(output_path)
        return output_path
