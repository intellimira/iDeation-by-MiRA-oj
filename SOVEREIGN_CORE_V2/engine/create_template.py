from pptx import Presentation
from pptx.util import Inches

def create_dummy_template():
    prs = Presentation()
    # Create a layout with placeholders
    # Index 1 is typically 'Title and Content'
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    
    title = slide.shapes.title
    title.text = "Master Title Placeholder"
    title.name = "Title 1"
    
    body = slide.placeholders[1]
    body.text = "Master Body Placeholder"
    body.name = "Text Placeholder 2"
    
    prs.save("sales_engine/templates/Brand_Master.pptx")
    print("Dummy template created at sales_engine/templates/Brand_Master.pptx")

if __name__ == "__main__":
    create_dummy_template()
